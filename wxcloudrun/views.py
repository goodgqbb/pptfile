from datetime import datetime
from flask import render_template, request,Flask
from run import app
from wxcloudrun.dao import delete_counterbyid, query_counterbyid, insert_counter, update_counterbyid
from wxcloudrun.model import Counters
from wxcloudrun.response import make_succ_empty_response, make_succ_response, make_err_response
import json
import requests
import numpy as np
import pandas as pd
import seaborn as sns
import lightgbm as lgb
from pptx import Presentation

#获取文件
def getfile(fileid):
    response = requests.get( 'https://api.weixin.qq.com/cgi-bin/token?grant_type=client_credential&appid=wxfb2997f507abf89e&secret=168726b557fb7221c96955cec59b3347',verify=False )
    access_token = response.json()['access_token']
    data ={
        "env": "prod-0gayxkvve034fe60",
        "file_list": [
        {
          "fileid":fileid,
          "max_age":86400
        }
        ]
      }
    #转json
    data = json.dumps(data, ensure_ascii=False).encode("utf-8")
    print(data)
    response = requests.post("https://api.weixin.qq.com/tcb/batchdownloadfile?access_token=" + access_token,data,verify=False)
    print(response.json())
    download_url = response.json()['file_list'][0]['download_url']
    filename =download_url.split("/")[-1]
    f = requests.get(download_url)
    path = "/app/wxcloudrun/"+filename
    #下载文件
    with open(path,"wb") as code:
         code.write(f.content)
    return path

@app.route('/', methods=['POST'])
def upload():
    fileid = request.json.get('fileid')
    path = getfile(fileid)
    data = []
    prs = Presentation(path)
    for slide in prs.slides: #遍历每页PPT
        for shape in slide.shapes: #遍历PPT中的每个形状
            if shape.has_text_frame: #判断该是否包含文本，保证有文本才提取
                for paragraph in shape.text_frame.paragraphs: #按文本框中的段落提取
                    data.append(paragraph.text) #提取一个段落的文本，就存到列表data中
    strr = "".join(data)
    ob = {
        "re":strr
    }
    c = json.dumps(ob)
    return c
